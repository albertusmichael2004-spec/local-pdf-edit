from __future__ import annotations

from io import BytesIO
from pathlib import Path

import fitz

from backend.core.errors import ConversionError
from backend.core.progress import report_fraction, report_progress

def pdf_to_pptx(input_path: Path, output_path: Path, dpi: int = 150) -> int:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError as exc:
        raise ConversionError("python-pptx is not installed. Run pip install -r requirements.txt.") from exc

    try:
        with fitz.open(input_path) as doc:
            if doc.needs_pass:
                raise ConversionError("Encrypted PDF. Unlock it before conversion.")
            if doc.page_count == 0:
                raise ConversionError("PDF has no pages.")
            first = doc[0].rect
            prs = Presentation()
            # Keep the PDF aspect ratio. PowerPoint slide dimensions are global.
            width_in = 10.0
            height_in = width_in * (first.height / first.width)
            prs.slide_width = Inches(width_in)
            prs.slide_height = Inches(height_in)
            blank = prs.slide_layouts[6]
            scale = dpi / 72.0
            for index, page in enumerate(doc, start=1):
                slide = prs.slides.add_slide(blank)
                pix = page.get_pixmap(matrix=fitz.Matrix(scale, scale), alpha=False)
                image_stream = BytesIO(pix.tobytes("png"))
                page_ratio = page.rect.width / page.rect.height
                slide_ratio = prs.slide_width / prs.slide_height
                if page_ratio >= slide_ratio:
                    pic_w = prs.slide_width
                    pic_h = int(pic_w / page_ratio)
                    left = 0
                    top = int((prs.slide_height - pic_h) / 2)
                else:
                    pic_h = prs.slide_height
                    pic_w = int(pic_h * page_ratio)
                    top = 0
                    left = int((prs.slide_width - pic_w) / 2)
                slide.shapes.add_picture(image_stream, left, top, width=pic_w, height=pic_h)
                report_fraction("Building PowerPoint slides", index, doc.page_count, 24, 90)
            # Remove the initial default slide if PowerPoint implementation created one.
            if len(prs.slides) > doc.page_count:
                slide_id = prs.slides._sldIdLst[0]
                prs.slides._sldIdLst.remove(slide_id)
            report_progress("Saving PowerPoint presentation", percent=94)
            prs.save(output_path)
            return doc.page_count
    except ConversionError:
        raise
    except Exception as exc:
        raise ConversionError(f"PDF to PowerPoint conversion failed: {exc}") from exc
