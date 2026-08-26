from __future__ import annotations

from io import BytesIO
from pathlib import Path

from backend.core.errors import ConversionError
from backend.services.shared.renderers.reportlab_common import reportlab_imports

def render_pptx_to_pdf(input_path: Path, output_path: Path) -> str:
    if input_path.suffix.lower() != ".pptx":
        raise ConversionError("Legacy .ppt files require LibreOffice. Save the file as .pptx or install LibreOffice.")
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as exc:
        raise ConversionError("python-pptx is not installed. Run pip install -r requirements.txt.") from exc
    rl = reportlab_imports()
    try:
        prs = Presentation(str(input_path))
        emu_per_inch = 914400
        width_pt = prs.slide_width / emu_per_inch * 72
        height_pt = prs.slide_height / emu_per_inch * 72
        c = rl["canvas"].Canvas(str(output_path), pagesize=(width_pt, height_pt))
        for slide in prs.slides:
            c.setFillColorRGB(1, 1, 1)
            c.rect(0, 0, width_pt, height_pt, fill=1, stroke=0)
            for shape in slide.shapes:
                x = shape.left / emu_per_inch * 72
                y_top = shape.top / emu_per_inch * 72
                w = shape.width / emu_per_inch * 72
                h = shape.height / emu_per_inch * 72
                y = height_pt - y_top - h
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        blob = shape.image.blob
                        c.drawImage(rl["ImageReader"](BytesIO(blob)), x, y, width=w, height=h, preserveAspectRatio=True, mask='auto')
                        continue
                except Exception:
                    pass
                if getattr(shape, "has_text_frame", False):
                    text = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text).strip()
                    if text:
                        c.setFillColorRGB(0.08, 0.08, 0.08)
                        font_size = max(8, min(28, h / max(1, len(text.splitlines())) * 0.42))
                        c.setFont("Helvetica", font_size)
                        text_obj = c.beginText(x, y + h - font_size)
                        text_obj.setLeading(font_size * 1.2)
                        for line in text.splitlines():
                            # ReportLab's base fonts are Latin-1; replace unsupported codepoints gracefully.
                            safe = line.encode("latin-1", "replace").decode("latin-1")
                            text_obj.textLine(safe[:250])
                        c.drawText(text_obj)
            c.showPage()
        c.save()
        return "Built-in PPTX renderer"
    except ConversionError:
        raise
    except Exception as exc:
        raise ConversionError(f"Built-in PPTX conversion failed: {exc}") from exc
